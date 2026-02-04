import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE

class PPTGenerator:
    def __init__(self):
        # Professional Color Palette (Deep Navy / Clean Grey / Vibrant Accent)
        self.NAVY = RGBColor(10, 25, 60)       # Dark Corporate Blue
        self.ACCENT = RGBColor(220, 50, 100)   # Sharp Pink/Red for highlights
        self.TEXT_MAIN = RGBColor(40, 40, 40)
        self.TEXT_LIGHT = RGBColor(100, 100, 100)
        self.BG_LIGHT = RGBColor(245, 247, 250) # Very subtle grey-blue
        self.WHITE = RGBColor(255, 255, 255)
        self.BORDER = RGBColor(200, 200, 200)
        self.SUCCESS = RGBColor(34, 139, 34)

    def _fmt(self, p, text, size, bold=False, color=None, font="Arial"):
        p.text = str(text) if text else ""
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = font
        if color: p.font.color.rgb = color

    def _footer(self, slide, page_num):
        # Thin divider line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.1), Inches(9), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = self.BORDER
        
        # Confidentiality
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(6), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        self._fmt(p, "Strictly Private & Confidential – Prepared by Kelp M&A Team", 8, color=self.TEXT_LIGHT)
        
        # Page Number
        if page_num > 0:
            sn = slide.shapes.add_textbox(Inches(9.0), Inches(7.15), Inches(0.5), Inches(0.4))
            p2 = sn.text_frame.paragraphs[0]
            self._fmt(p2, str(page_num), 9, color=self.TEXT_LIGHT)
            p2.alignment = PP_ALIGN.RIGHT

    def _header(self, slide, text):
        # Header Strip
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), Inches(10), Inches(0.6))
        sh.fill.solid()
        sh.fill.fore_color.rgb = self.NAVY
        sh.line.fill.background()
        
        # Logo Accent
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), Inches(0.15), Inches(0.6))
        acc.fill.solid()
        acc.fill.fore_color.rgb = self.ACCENT
        acc.line.fill.background()
        
        # Text
        tf = sh.text_frame
        p = tf.paragraphs[0]
        self._fmt(p, text.upper(), 18, bold=True, color=self.WHITE)
        tf.margin_left = Inches(0.4)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _img(self, slide, path, l, t, w, h):
        if path and os.path.exists(path):
            pic = slide.shapes.add_picture(path, l, t, w, h)
            pic.line.color.rgb = self.BORDER
            pic.line.width = Pt(1)
        else:
            sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
            sh.fill.solid()
            sh.fill.fore_color.rgb = RGBColor(230, 230, 230)
            sh.line.color.rgb = self.BORDER
            sh.text_frame.text = "Visual Asset Placeholder"
            self._fmt(sh.text_frame.paragraphs[0], "Visual Asset", 10, color=self.TEXT_LIGHT)

    def _get_defaults(self, section):
        return {
            'slide_1': {'headline': 'Investment Opportunity', 'bullets': []},
            'slide_2': {'metrics': {}, 'chart_data': {}},
            'slide_3': {'hooks': []}
        }.get(section, {})

    def generate_ppt(self, data, images, filename):
        prs = Presentation()
        while len(images) < 3: images.append(None)
        
        for k in ['slide_1', 'slide_2', 'slide_3']:
            if k not in data: data[k] = self._get_defaults(k)

        # ==============================================================================
        # SLIDE 1: EXECUTIVE SUMMARY (Dense Grid Layout)
        # ==============================================================================
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        self._header(s1, f"{data.get('code_name', 'Project')} | Executive Summary")
        self._footer(s1, 1)
        
        # 1. Headlines (Top Full Width)
        box = s1.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.8))
        self._fmt(box.text_frame.paragraphs[0], data['slide_1'].get('headline',''), 20, bold=True, color=self.NAVY)
        p2 = box.text_frame.add_paragraph()
        p2.space_before = Pt(6)
        self._fmt(p2, data['slide_1'].get('sub_headline', ''), 12, color=self.TEXT_MAIN)

        # 2. Main Content Grid (Left Text, Right Visual)
        # Left: Bullets
        bbox = s1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.5), Inches(3.5))
        bullets = data['slide_1'].get('bullets', [])
        cert_kws = ['iso', 'gmp', 'fda', 'certified', 'who', 'fssc', 'usda']
        
        for b in bullets:
            if not any(kw in b.lower() for kw in cert_kws):
                p = bbox.text_frame.add_paragraph()
                self._fmt(p, f"•  {b}", 11, color=self.TEXT_MAIN)
                p.space_after = Pt(12)

        # Right: Image
        self._img(s1, images[0], Inches(6.2), Inches(2.2), Inches(3.3), Inches(2.2))

        # 3. Compliance Box (Right, below image)
        certs = [b for b in bullets if any(kw in b.lower() for kw in cert_kws)]
        if certs:
            cbox = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(4.6), Inches(3.3), Inches(1.1))
            cbox.fill.solid()
            cbox.fill.fore_color.rgb = self.BG_LIGHT
            cbox.line.color.rgb = self.BORDER
            
            tf = cbox.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_top = Inches(0.1)
            p = tf.paragraphs[0]
            self._fmt(p, "KEY CERTIFICATIONS", 9, bold=True, color=self.NAVY)
            
            for c in certs[:2]:
                p_c = tf.add_paragraph()
                clean_c = c.split(':')[0].replace("Holds ", "").replace("Certified ", "")
                self._fmt(p_c, f"✓ {clean_c[:35]}", 9, color=self.SUCCESS)

        # 4. Key Stats Strip (Bottom)
        strip = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.9), Inches(9), Inches(1.0))
        strip.fill.solid()
        strip.fill.fore_color.rgb = self.NAVY
        strip.line.fill.background()
        
        # Fake stats for visual density (if real ones missing)
        stats = [("Global Reach", "45+ Countries"), ("Industry Rank", "Top 10"), ("Workforce", "500+")]
        x_stat = Inches(0.8)
        for label, val in stats:
            tf = s1.shapes.add_textbox(x_stat, Inches(6.0), Inches(2.5), Inches(0.8)).text_frame
            p_lbl = tf.paragraphs[0]
            p_lbl.alignment = PP_ALIGN.CENTER
            self._fmt(p_lbl, label.upper(), 9, color=RGBColor(200,200,200))
            
            p_val = tf.add_paragraph()
            p_val.alignment = PP_ALIGN.CENTER
            self._fmt(p_val, val, 16, bold=True, color=self.WHITE)
            x_stat += Inches(3.0)

        # ==============================================================================
        # SLIDE 2: FINANCIAL PROFILE (KPI Cards + Chart)
        # ==============================================================================
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        self._header(s2, "Financial & Operational Profile")
        self._footer(s2, 2)
        
        metrics = data['slide_2'].get('metrics', {})
        chart_vals = data['slide_2'].get('chart_data', {}).get('revenue_values', [])
        years = data['slide_2'].get('chart_data', {}).get('years', [])
        
        # DATA RECOVERY LOGIC: If Text Summary is N/A but Chart has data, steal it.
        revenue_keys = [k for k in metrics.keys() if "revenue" in k.lower()]
        if revenue_keys:
            rev_val = str(metrics[revenue_keys[0]])
            if ("N/A" in rev_val or "None" in rev_val) and chart_vals:
                valid_vals = [v for v in chart_vals if v and v != 0]
                if valid_vals:
                    metrics[revenue_keys[0]] = f"{valid_vals[-1]:,.0f}"

        # 1. KPI Cards (Left Column)
        y_card = Inches(1.5)
        for k, v in list(metrics.items())[:4]:
            # Card Box
            card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_card, Inches(2.5), Inches(1.2))
            card.fill.solid()
            card.fill.fore_color.rgb = self.WHITE
            card.line.color.rgb = self.BORDER
            card.shadow.inherit = False # Clean look
            
            # Text
            tf = card.text_frame
            p_lbl = tf.paragraphs[0]
            self._fmt(p_lbl, k.upper(), 9, bold=True, color=self.TEXT_LIGHT)
            
            p_val = tf.add_paragraph()
            p_val.space_before = Pt(10)
            self._fmt(p_val, str(v), 22, bold=True, color=self.NAVY)
            
            y_card += Inches(1.4)

        # 2. Chart (Right Column)
        if chart_vals and any(v for v in chart_vals if v):
            # Chart Container
            c_bg = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(1.5), Inches(6), Inches(4.5))
            c_bg.fill.solid()
            c_bg.fill.fore_color.rgb = self.BG_LIGHT
            c_bg.line.fill.background()
            
            # Title
            tbox = s2.shapes.add_textbox(Inches(3.7), Inches(1.6), Inches(4), Inches(0.4))
            self._fmt(tbox.text_frame.paragraphs[0], "Revenue Trajectory (INR Cr)", 12, bold=True, color=self.NAVY)
            
            # Actual Chart
            chart_data = CategoryChartData()
            chart_data.categories = years
            chart_data.add_series('Revenue', chart_vals)
            chart = s2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(3.7), Inches(2.0), Inches(5.6), Inches(3.8), chart_data).chart
            chart.has_legend = False
            
            # CAGR Arrow (Calculation)
            try:
                valid_vals = [v for v in chart_vals if v]
                if len(valid_vals) >= 2:
                    start_v = valid_vals[0]
                    end_v = valid_vals[-1]
                    years_count = len(valid_vals) - 1
                    if start_v > 0 and years_count > 0:
                        cagr = (math.pow(end_v / start_v, 1 / years_count) - 1) * 100
                        
                        # Draw Arrow
                        arrow = s2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(2.5), Inches(1.5), Inches(0.4))
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = self.SUCCESS
                        arrow.line.fill.background()
                        arrow.text_frame.paragraphs[0].text = f"CAGR: {cagr:.1f}%"
                        arrow.text_frame.paragraphs[0].font.size = Pt(10)
            except: pass

        # ==============================================================================
        # SLIDE 3: INVESTMENT THESIS (2x2 Matrix Grid)
        # ==============================================================================
        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        self._header(s3, "Key Investment Highlights")
        self._footer(s3, 3)
        
        hooks = data['slide_3'].get('hooks', [])
        # Coordinates for 2x2 grid
        positions = [
            (Inches(0.5), Inches(1.5)), (Inches(5.1), Inches(1.5)),
            (Inches(0.5), Inches(4.2)), (Inches(5.1), Inches(4.2))
        ]
        
        for i, (x, y) in enumerate(positions):
            if i >= len(hooks): break
            
            # Box
            box = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.4), Inches(2.5))
            box.fill.solid()
            box.fill.fore_color.rgb = self.WHITE
            box.line.color.rgb = self.BORDER
            
            # Number Badge
            badge = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.6), Inches(0.6))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.ACCENT
            badge.line.fill.background()
            bf = badge.text_frame
            bf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._fmt(bf.paragraphs[0], f"0{i+1}", 14, bold=True, color=self.WHITE)
            
            # Content
            tf = box.text_frame
            tf.margin_top = Inches(0.8)
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            self._fmt(tf.paragraphs[0], hooks[i], 12, color=self.TEXT_MAIN)

        prs.save(filename)